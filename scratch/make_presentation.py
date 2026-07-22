import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5) # 16:9 widescreen layout

# Custom Luxury Palette
COLOR_DARK = RGBColor(44, 38, 35)      # #2C2623 - Charcoal
COLOR_LIGHT = RGBColor(252, 250, 246)  # #FCFAF6 - Cream
COLOR_GOLD = RGBColor(197, 168, 128)   # #C5A880 - Gold
COLOR_MUTED = RGBColor(110, 99, 93)    # #6E635D - Muted Taupe
COLOR_WHITE = RGBColor(255, 255, 255)

# Target folder
screenshot_dir = r"c:\Users\Yash Bhatt\Desktop\Demo\Screeshots"

# Regex to parse the timestamps from WhatsApp image filenames
def get_image_sort_key(filename):
    # Match: WhatsApp Image 2026-07-22 at 1.02.08 AM.jpeg
    # Extra matches for optional parenthesis: (1), (2)
    match = re.search(r'(\d+)\.(\d+)\.(\d+)\s+(AM|PM)(?:\s+\((\d+)\))?', filename)
    if match:
        hour, minute, second, meridiem, suffix = match.groups()
        hour = int(hour)
        minute = int(minute)
        second = int(second)
        
        # Convert to 24h
        if meridiem == 'PM' and hour != 12:
            hour += 12
        elif meridiem == 'AM' and hour == 12:
            hour = 0
            
        suffix = int(suffix) if suffix else 0
        return (hour, minute, second, suffix)
    return (0, 0, 0, 0)

# Retrieve and sort images
all_files = [f for f in os.listdir(screenshot_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
sorted_images = sorted(all_files, key=get_image_sort_key)

print(f"Found {len(sorted_images)} images. Generating presentation...")

# Helper to apply background color to a slide
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Helper to add standard text
def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=COLOR_DARK, font_name="Georgia", align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

# 1. Slide 1: Welcome Title Slide (Dark Theme)
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1, COLOR_DARK)

add_textbox(slide1, Inches(1), Inches(2.2), Inches(11.3), Inches(1), 
            "EVENTSPHERE PORTAL", font_size=20, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5), 
            "End-to-End Automated Testing Report", font_size=42, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(5.0), Inches(11.3), Inches(1), 
            "Cypress E2E Testing • Selenium Webdriver Scripts • Selenium IDE Suite", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# Mapping of images to sections and steps
# We have 25 images total. We will divide them into 3 logical chapters.
image_steps = [
    # Section 1: Selenium IDE
    {"sec": "Selenium IDE", "title": "Dashboard Recorder Setup", "desc": "Launching Selenium IDE Chrome extension and initializing a new test suite project for EventSphere."},
    {"sec": "Selenium IDE", "title": "Base URL Configuration", "desc": "Configuring target site URL parameters to bind IDE playback to the live Render host address."},
    {"sec": "Selenium IDE", "title": "Login Flow Recording", "desc": "Recording interactive user click actions on input boxes, email fields, and password submissions."},
    {"sec": "Selenium IDE", "title": "DOM Element Selectors", "desc": "Verifying target IDs, class structures, and visual elements to guarantee assertion success during playback."},
    {"sec": "Selenium IDE", "title": "Assertion Command Setup", "desc": "Adding custom assertions to verify that dashboard headers like 'Active Passes' render correctly."},
    {"sec": "Selenium IDE", "title": "Playback Execution", "desc": "Executing the recorded test suite step-by-step. Selenium IDE drives the active browser session."},
    {"sec": "Selenium IDE", "title": "Interactive Flow Test", "desc": "Navigating between public events explorer, event details view, and user portal pages during replay."},
    {"sec": "Selenium IDE", "title": "IDE Playback Log Verification", "desc": "Confirming all command sequences completed successfully with zero error indicators."},

    # Section 2: Selenium Webdriver Terminal Execution
    {"sec": "Selenium WebDriver", "title": "Terminal Script Configuration", "desc": "Setting up the Node.js automation script using selenium-webdriver bindings to drive Chrome."},
    {"sec": "Selenium WebDriver", "title": "Test Credentials Bind", "desc": "Adding secure credential hooks and target environment URL declarations inside the selenium test script."},
    {"sec": "Selenium WebDriver", "title": "Chrome Driver Initializing", "desc": "Executing the script via 'node selenium_test.js' and instantiating the Chrome browser instance."},
    {"sec": "Selenium WebDriver", "title": "Automated Page Navigation", "desc": "Webdriver navigating Chrome to the live EventSphere login page at the specified Render address."},
    {"sec": "Selenium WebDriver", "title": "Automated Form Input", "desc": "Selenium Webdriver locating email/password inputs and writing test credentials dynamically."},
    {"sec": "Selenium WebDriver", "title": "Submit Button Action", "desc": "Simulating the button click to submit authorization parameters to the backend server."},
    {"sec": "Selenium WebDriver", "title": "Wait & Assert Redirect", "desc": "Waiting for URL routing to contain '/dashboard' to confirm authorization connection succeeded."},
    {"sec": "Selenium WebDriver", "title": "Terminal Log Success", "desc": "Verifying that the test script exits cleanly and prints 'Test Passed' inside the PowerShell console."},

    # Section 3: Cypress E2E Testing Suite
    {"sec": "Cypress E2E", "title": "Cypress Project Launch", "desc": "Launching Cypress test runner and configuring specs structure for frontend application testing."},
    {"sec": "Cypress E2E", "title": "spec.cy.js Declaration", "desc": "Writing pure JavaScript Cypress assertions to verify main pages, forms, and login flows."},
    {"sec": "Cypress E2E", "title": "Homepage Visual Checks", "desc": "Cypress visiting the homepage and checking branding elements, typography, and Call-to-Action states."},
    {"sec": "Cypress E2E", "title": "About Page Assertion", "desc": "Navigating to the About page and verifying that philosophy paragraphs and headings are correctly visible."},
    {"sec": "Cypress E2E", "title": "Contact Concierge Form", "desc": "Testing the concierge contact form by filling name, email, details, and clicking the submit button."},
    {"sec": "Cypress E2E", "title": "Toast Feedback Assertion", "desc": "Asserting that the premium success toast message appears correctly after submitting a contact message."},
    {"sec": "Cypress E2E", "title": "Search Bar Input Filter", "desc": "Navigating to Events catalog and typing inside the filter search bar to verify responsiveness of listing."},
    {"sec": "Cypress E2E", "title": "Attendee Login Test", "desc": "Executing automated login using Cy inputs and asserting dashboard redirect is resolved instantly."},
    {"sec": "Cypress E2E", "title": "Dashboard Navigation Verification", "desc": "Clicking dashboard tabs (My Tickets, Wishlist, Settings) to verify client-side router redirects."}
]

# Walk through each sorted image and build slides
for idx, img_name in enumerate(sorted_images):
    # Fallback details if index exceeds step details
    if idx < len(image_steps):
        step_info = image_steps[idx]
    else:
        step_info = {"sec": "Automation Step", "title": f"Operation View {idx + 1}", "desc": "Visual verification screen capture illustrating testing sequence progression."}

    # Add content slide (Light Theme)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_LIGHT)

    # 1. Category Tag (Top-left)
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4), 
                step_info["sec"].upper(), font_size=10, bold=True, color=COLOR_GOLD, font_name="Arial")

    # 2. Main Title (Top-left)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(5.0), Inches(0.8), 
                f"{idx + 1}. {step_info['title']}", font_size=26, bold=True, color=COLOR_DARK, font_name="Georgia")

    # 3. Description (Middle-left)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(4.5), Inches(3.5), 
                step_info["desc"], font_size=14, color=COLOR_MUTED, font_name="Arial")
    
    # 4. Slide Index Footer (Bottom-left)
    add_textbox(slide, Inches(0.8), Inches(6.5), Inches(3.0), Inches(0.4), 
                f"Slide {idx + 2} of {len(sorted_images) + 1}", font_size=9, color=COLOR_GOLD)

    # 5. Insert Image on the Right
    img_path = os.path.join(screenshot_dir, img_name)
    try:
        # Standard widescreen coordinates for image:
        # Left = 5.8", Top = 0.9", Width = 6.8", Height = 5.5" (keeps aspect ratio)
        slide.shapes.add_picture(img_path, Inches(5.8), Inches(0.9), width=Inches(6.8), height=Inches(5.5))
    except Exception as e:
        print(f"Error loading image {img_name}: {e}")

# Save presentation
output_pptx = r"c:\Users\Yash Bhatt\Desktop\Demo\EventSphere_Testing_Report.pptx"
prs.save(output_pptx)
print(f"PowerPoint Presentation saved successfully: {output_pptx}")
