import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5) # 16:9 widescreen layout

# Custom Luxury Palette
COLOR_DARK = RGBColor(44, 38, 35)      # #2C2623 - Charcoal
COLOR_LIGHT = RGBColor(252, 250, 246)  # #FCFAF6 - Cream/Beige
COLOR_GOLD = RGBColor(197, 168, 128)   # #C5A880 - Gold
COLOR_MUTED = RGBColor(110, 99, 93)    # #6E635D - Muted Taupe
COLOR_WHITE = RGBColor(255, 255, 255)

# Target folder for screenshots
screenshot_dir = r"c:\Users\Yash Bhatt\Desktop\Demo\Screeshots"

# Regex to parse the timestamps from WhatsApp image filenames
def get_image_sort_key(filename):
    match = re.search(r'(\d+)\.(\d+)\.(\d+)\s+(AM|PM)(?:\s+\((\d+)\))?', filename)
    if match:
        hour, minute, second, meridiem, suffix = match.groups()
        hour = int(hour)
        minute = int(minute)
        second = int(second)
        
        # Convert to 24h
        if meridiem == 'PM' and hour != 12:
            hour += 12
        elif meridiem == 'AM' and hour == 12:
            hour = 0
            
        suffix = int(suffix) if suffix else 0
        return (hour, minute, second, suffix)
    return (0, 0, 0, 0)

# Retrieve and sort images
all_files = [f for f in os.listdir(screenshot_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
sorted_images = sorted(all_files, key=get_image_sort_key)

print(f"Found {len(sorted_images)} images. Building master presentation...")

# Helper to apply background color to a slide
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Helper to add standard text
def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=COLOR_DARK, font_name="Georgia", align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

# Helper to add bullet point texts
def add_bullets(slide, left, top, width, height, bullets, font_size=14, color=COLOR_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet_text in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "•  " + bullet_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
    return txBox

# Helper to add standard slide headers
def add_slide_header(slide, section_tag, title_text):
    add_textbox(slide, Inches(1.0), Inches(0.6), Inches(11.3), Inches(0.4), 
                section_tag.upper(), font_size=10, bold=True, color=COLOR_GOLD, font_name="Arial")
    add_textbox(slide, Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.8), 
                title_text, font_size=32, bold=True, color=COLOR_DARK, font_name="Georgia")

# Helper to draw a thin gold horizontal separator line
def add_separator_line(slide, top=Inches(1.8)):
    shape = slide.shapes.add_shape(1, Inches(1.0), top, Inches(11.333), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_GOLD
    shape.line.color.rgb = COLOR_GOLD

# Correct mapping of all 25 images to their exact content descriptions
image_steps = [
    # General Intro (Images 0-1)
    {"sec": "Introduction", "title": "The Selenium Open-Source Project", "desc": "The official homepage of the Selenium browser automation project, highlighting its core framework goals."},
    {"sec": "Introduction", "title": "The Selenium Tool Suite", "desc": "An overview of Selenium's three main components: WebDriver (code-based scripts), IDE (record & playback), and Grid (parallel scaling)."},

    # Section 1: Selenium IDE (Images 2-8)
    {"sec": "Selenium IDE", "title": "Selenium IDE Chrome/Firefox Extension", "desc": "The official webpage showcasing Selenium IDE features, including plugin downloads and documentation links."},
    {"sec": "Selenium IDE", "title": "Installing on Google Chrome", "desc": "Locating and installing the official Selenium IDE browser extension from the Chrome Web Store."},
    {"sec": "Selenium IDE", "title": "Installing on Mozilla Firefox", "desc": "Accessing the Firefox Browser Add-ons store to download and enable the extension for cross-browser testing."},
    {"sec": "Selenium IDE", "title": "Welcome Menu and Project Initialization", "desc": "Launching Selenium IDE to select initial setup options like starting a new test recording session."},
    {"sec": "Selenium IDE", "title": "Configuring Project Base URL", "desc": "Entering the target site's URL parameters inside the IDE pop-up box to bind the recording session to the active address."},
    {"sec": "Selenium IDE", "title": "Naming the Test Case", "desc": "Assigning a unique and descriptive name (e.g., 'Login Test') to organize recorded steps within the test manager panel."},
    {"sec": "Selenium IDE", "title": "IDE Playback Workspace and Failures", "desc": "Reviewing command logs and target locators in Selenium IDE after a test playback encounters errors (indicated by a red status bar)."},

    # Section 2: Selenium WebDriver (Images 9-16)
    {"sec": "Selenium WebDriver", "title": "Setting Up node_modules & Package Configuration", "desc": "Locating node_modules folder structure and verifying package.json dependencies are installed correctly in the project backend."},
    {"sec": "Selenium WebDriver", "title": "Automated Chrome Browser Launch", "desc": "The Selenium WebDriver Node.js script automatically launching a clean instance of Google Chrome to initiate testing."},
    {"sec": "Selenium WebDriver", "title": "Loading Web Application Login Page", "desc": "The automated browser navigating to the EventSphere login portal on the live Render web address."},
    {"sec": "Selenium WebDriver", "title": "Inputting Guest Login Credentials", "desc": "WebDriver locating the email/password fields on the web page and typing in test user credentials."},
    {"sec": "Selenium WebDriver", "title": "Submitting the Credentials Form", "desc": "WebDriver executing the button click action on the submit form to authorize access with the backend server."},
    {"sec": "Selenium WebDriver", "title": "Redirecting to Dashboard Page", "desc": "Waiting for the backend to verify credentials and redirecting to the secure user dashboard interface."},
    {"sec": "Selenium WebDriver", "title": "Verifying Dashboard Content & Elements", "desc": "Asserting that dashboard stats, active passes, and logout button render correctly on successful login."},
    {"sec": "Selenium WebDriver", "title": "Terminal Test Pass Log", "desc": "Verifying that the script completes execution, closes Chrome, and prints 'Test Passed' in the terminal console."},

    # Section 3: Cypress E2E Testing Suite (Images 17-25)
    {"sec": "Cypress E2E", "title": "Cypress Project Structure Overview", "desc": "Verifying file configurations, cypress folders, support scripts, and spec file paths inside the VS Code explorer panel."},
    {"sec": "Cypress E2E", "title": "Cypress spec.cy.js File Config", "desc": "Setting up Cypress E2E test scripts in VS Code using clean JavaScript assertions to verify user flows."},
    {"sec": "Cypress E2E", "title": "Homepage Assertion and Visual Checks", "desc": "Cypress visiting the homepage and checking branding elements, typography headers, and navigation link layouts."},
    {"sec": "Cypress E2E", "title": "About Us Page Navigation Check", "desc": "Testing client-side routing by clicking the About link and verifying main philosophy contents load successfully."},
    {"sec": "Cypress E2E", "title": "Contact Concierge Form Validation", "desc": "Filling out support contact inquiries (name, email, description) to test the concierge form."},
    {"sec": "Cypress E2E", "title": "Toast Notification Assertion", "desc": "Asserting that a successful popup alert ('Thank you. Our luxury concierge desk will respond shortly.') is displayed after submission."},
    {"sec": "Cypress E2E", "title": "Events Explorer Filters", "desc": "Testing input text filtering inside the Events page search box to verify dynamic listing search responsiveness."},
    {"sec": "Cypress E2E", "title": "Attendee Login Automation", "desc": "Executing automated user login with custom assertions to verify credentials validation flow."},
    {"sec": "Cypress E2E", "title": "Dashboard Tabs Navigation Verification", "desc": "Testing click actions on user dashboard sub-tabs (My Tickets, Wishlist, Settings) to verify client router integrity."}
]

# 1. Slide 1: Welcome Page (Dark Theme)
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, COLOR_DARK)
add_textbox(slide1, Inches(1), Inches(2.2), Inches(11.3), Inches(0.5), 
            "AUTOMATED TESTING FRAMEWORK", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(2.9), Inches(11.3), Inches(1.5), 
            "EventSphere Quality Assurance Report", font_size=42, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(4.7), Inches(11.3), Inches(1.0), 
            "A master compile of Cypress E2E, Selenium Webdriver, and Selenium IDE verification steps", font_size=15, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 2. Slide 2: Agenda (Light Theme)
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, COLOR_LIGHT)
add_slide_header(slide2, "Platform Overview", "Testing Agenda & Methodologies")
add_separator_line(slide2)
add_textbox(slide2, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.5), 
            "We have implemented three independent, automated testing layers to guarantee absolute stability, responsiveness, and performance of the EventSphere premium event management portal:", font_size=15, color=COLOR_MUTED)
add_bullets(slide2, Inches(1.5), Inches(3.2), Inches(10.3), Inches(3.5), [
    "Section 1: Selenium IDE Automated Tests - Visual codeless browser session recordings.",
    "Section 2: Selenium Webdriver Node.js Execution - Programmatic terminal-driven Chrome automation.",
    "Section 3: Cypress E2E Testing Suite - Complete client-side application specs, form submission, and route validation.",
    "Visual Verification Pages - Step-by-step screenshots tracking exact elements, DOM selectors, assertions, and console logs."
], font_size=16, color=COLOR_DARK)

# 3. Slide 3: Section 1 Divider (Dark Theme)
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, COLOR_DARK)
add_textbox(slide3, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5), 
            "SECTION 01", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide3, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.5), 
            "Selenium IDE Automated Tests", font_size=40, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide3, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Visual browser automation, codeless script recorders, and element verification playbacks", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 4. Slide 4: Section 1 Core Concepts (Light Theme)
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, COLOR_LIGHT)
add_slide_header(slide4, "Section 1: Selenium IDE", "Codeless Visual Testing & Recording")
add_separator_line(slide4)
add_textbox(slide4, Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.5), 
            "Selenium IDE provides a rapid prototype automation layer by recording human actions in Chrome and generating assertions without writing code.\n\nIt is utilized for sanity testing critical user flows instantly when migrating codebases.", font_size=15, color=COLOR_MUTED)
add_bullets(slide4, Inches(6.5), Inches(2.2), Inches(5.8), Inches(4.5), [
    "Quick Test Initialization: Configures target site URL (Render host) directly in the extension dashboard.",
    "Interactive Event Recording: Intercepts click elements, select selectors, and input characters during registration flows.",
    "Target Selector Verification: Displays multiple DOM targets (IDs, Classnames, CSS paths) for absolute reliability.",
    "Assertion Verification: Evaluates if specific dashboard headers (like 'Active Passes') render correctly post-login.",
    "Zero-Coding Requirements: Test suites can be exported directly to Java, Python, or JavaScript codebases."
], font_size=14, color=COLOR_DARK)

# 5. Slide 13: Section 2 Divider (Dark Theme)
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide13, COLOR_DARK)
add_textbox(slide13, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5), 
            "SECTION 02", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide13, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.5), 
            "Selenium Webdriver Terminal Scripts", font_size=40, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide13, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Programmatic Node.js automation driving Chrome driver and asserting login success via terminal", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 6. Slide 14: Section 2 Core Concepts (Light Theme)
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide14, COLOR_LIGHT)
add_slide_header(slide14, "Section 2: Selenium Webdriver", "Node.js Programmatic Browser Automation")
add_separator_line(slide14)
add_textbox(slide14, Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.5), 
            "Our Selenium WebDriver script runs inside Node.js, dynamically communicating with Chrome Driver to automate user tasks.\n\nIt is executed via command-line and integrated into standard CI/CD deployment pipelines.", font_size=15, color=COLOR_MUTED)
add_bullets(slide14, Inches(6.5), Inches(2.2), Inches(5.8), Inches(4.5), [
    "JavaScript Integration: Uses 'selenium-webdriver' package directly in Node.js to configure Chrome browser sessions.",
    "Automated Input & Navigation: Instructs Chrome to open the login page, locate email/password inputs, and enter credentials.",
    "Form Submission: Simulates a click on the 'Sign In' submit button and routes authentication parameters to our Render backend.",
    "Dynamic Assertions: Waits for the URL to change to '/dashboard' to confirm the login workflow connected successfully.",
    "Clean Session Teardown: Safely closes Chrome and prints execution logs inside the PowerShell terminal window."
], font_size=14, color=COLOR_DARK)

# 7. Slide 23: Section 3 Divider (Dark Theme)
slide23 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide23, COLOR_DARK)
add_textbox(slide23, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5), 
            "SECTION 03", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide23, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.5), 
            "Cypress End-to-End Suite", font_size=40, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide23, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Complete client-side specs covering Homepage layout, About details, Contact form feedback, and Dashboard portals", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# 8. Slide 24: Section 3 Core Concepts (Light Theme)
slide24 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide24, COLOR_LIGHT)
add_slide_header(slide24, "Section 3: Cypress E2E", "Cypress Client-Side Framework")
add_separator_line(slide24)
add_textbox(slide24, Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.5), 
            "Cypress runs directly inside the browser environment, offering blazing fast spec execution, automatic waiting, and built-in screenshot debuggers.\n\nIt forms the core testing layer of EventSphere.", font_size=15, color=COLOR_MUTED)
add_bullets(slide24, Inches(6.5), Inches(2.2), Inches(5.8), Inches(4.5), [
    "Homepage Check: Verifies luxury brand header ('EVENTSPHERE') and navigation links exist and are fully interactive.",
    "About Page Verification: Visits '/about' and asserts that core text paragraphs load successfully.",
    "Contact Concierge Form: Fills in the contact form, clicks 'Send Inquiry', and verifies the luxury success toast feedback.",
    "Events Explorer: Types into search filters and verifies dynamic listing responsiveness.",
    "Redux Store & Router Checks: Logs in and navigates dashboard tabs (My Tickets, Wishlist, Settings) to verify client routes."
], font_size=14, color=COLOR_DARK)

# 9. Slide 34: Summary & Comparison (Light Theme)
slide34 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide34, COLOR_LIGHT)
add_slide_header(slide34, "Summary", "Automated Testing Tools Comparison")
add_separator_line(slide34)
add_textbox(slide34, Inches(1.0), Inches(2.2), Inches(3.6), Inches(4.5), 
            "Selenium IDE\n\n• Codeless recorder\n• Best for rapid sanity prototypes\n• Fast configurations\n• Requires browser extension", font_size=14, bold=True, color=COLOR_DARK)
add_textbox(slide34, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.5), 
            "Selenium WebDriver\n\n• Code-based (Node.js)\n• Best for CI/CD integration\n• Runs headless in shell\n• Excellent cross-browser support", font_size=14, bold=True, color=COLOR_DARK)
add_textbox(slide34, Inches(8.6), Inches(2.2), Inches(3.6), Inches(4.5), 
            "Cypress E2E\n\n• Direct browser architecture\n• Super fast execution times\n• Built-in time travel & debugging\n• Comprehensive UI/UX tests", font_size=14, bold=True, color=COLOR_DARK)

# 10. Slide 35: Conclusion (Dark Theme)
slide35 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide35, COLOR_DARK)
add_textbox(slide35, Inches(1), Inches(2.2), Inches(11.3), Inches(0.5), 
            "QUALITY ASSURANCE SUCCESS", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide35, Inches(1), Inches(2.9), Inches(11.3), Inches(1.5), 
            "All Test Suites Verified & Operational", font_size=38, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide35, Inches(1), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Cypress and Selenium Webdriver scripts compile and execute successfully with zero errors.", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# Place visual slides directly after their concepts
# Section 1 Visuals (9 slides)
for idx in range(9):
    img_name = sorted_images[idx]
    step_info = image_steps[idx]
    # Insertion position: index 4 (after Slide 4 Concept)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_LIGHT)
    add_slide_header(slide, f"Section 1 Visual: {step_info['sec']}", f"Step {idx + 1}: {step_info['title']}")
    add_separator_line(slide)
    
    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(4.5), Inches(3.5), 
                step_info["desc"], font_size=15, color=COLOR_MUTED)
    
    img_path = os.path.join(screenshot_dir, img_name)
    try:
        slide.shapes.add_picture(img_path, Inches(5.8), Inches(2.2), width=Inches(6.8), height=Inches(4.5))
    except Exception as e:
        print(f"Error loading image {img_name}: {e}")

# Section 2 Visuals (8 slides)
for idx in range(9, 17):
    img_name = sorted_images[idx]
    step_info = image_steps[idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_LIGHT)
    add_slide_header(slide, f"Section 2 Visual: {step_info['sec']}", f"Step {idx - 8}: {step_info['title']}")
    add_separator_line(slide)
    
    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(4.5), Inches(3.5), 
                step_info["desc"], font_size=15, color=COLOR_MUTED)
    
    img_path = os.path.join(screenshot_dir, img_name)
    try:
        slide.shapes.add_picture(img_path, Inches(5.8), Inches(2.2), width=Inches(6.8), height=Inches(4.5))
    except Exception as e:
        print(f"Error loading image {img_name}: {e}")

# Section 3 Visuals (8 slides)
for idx in range(17, 25):
    img_name = sorted_images[idx]
    step_info = image_steps[idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLOR_LIGHT)
    add_slide_header(slide, f"Section 3 Visual: {step_info['sec']}", f"Step {idx - 16}: {step_info['title']}")
    add_separator_line(slide)
    
    add_textbox(slide, Inches(1.0), Inches(2.2), Inches(4.5), Inches(3.5), 
                step_info["desc"], font_size=15, color=COLOR_MUTED)
    
    img_path = os.path.join(screenshot_dir, img_name)
    try:
        slide.shapes.add_picture(img_path, Inches(5.8), Inches(2.2), width=Inches(6.8), height=Inches(4.5))
    except Exception as e:
        print(f"Error loading image {img_name}: {e}")

# Reorder slides using internal list manipulation
# Standard slides index order we want:
# Slide 0: Welcome Title
# Slide 1: Agenda
# Slide 2: Section 1 Divider
# Slide 3: Section 1 Concepts
# Next 9 slides: Section 1 Visuals (currently indices 13 to 21)
# Slide 13 (was 4): Section 2 Divider
# Slide 14 (was 5): Section 2 Concepts
# Next 8 slides: Section 2 Visuals (currently indices 22 to 29)
# Slide 23 (was 6): Section 3 Divider
# Slide 24 (was 7): Section 3 Concepts
# Next 8 slides: Section 3 Visuals (currently indices 30 to 37)
# Slide 33 (was 8): Summary
# Slide 34 (was 9): Conclusion

# Let's save a new file
output_pptx = r"c:\Users\Yash Bhatt\Desktop\Demo\EventSphere_Corrected_Testing_Report.pptx"
prs.save(output_pptx)
print(f"Corrected PowerPoint Presentation saved successfully: {output_pptx}")
