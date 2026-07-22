from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5) # 16:9 widescreen layout

# Custom Luxury Palette (EventSphere Branding)
COLOR_DARK = RGBColor(44, 38, 35)      # #2C2623 - Charcoal
COLOR_LIGHT = RGBColor(252, 250, 246)  # #FCFAF6 - Cream/Beige
COLOR_GOLD = RGBColor(197, 168, 128)   # #C5A880 - Gold
COLOR_MUTED = RGBColor(110, 99, 93)    # #6E635D - Muted Taupe
COLOR_WHITE = RGBColor(255, 255, 255)

# Helper to apply background color to a slide
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Helper to add text boxes
def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=COLOR_DARK, font_name="Georgia", align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

# Helper to add bullet point texts
def add_bullets(slide, left, top, width, height, bullets, font_size=14, color=COLOR_DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for idx, bullet_text in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "•  " + bullet_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
    return txBox

# Helper to add standard slide headers
def add_slide_header(slide, section_tag, title_text):
    add_textbox(slide, Inches(1.0), Inches(0.6), Inches(11.3), Inches(0.4), 
                section_tag.upper(), font_size=10, bold=True, color=COLOR_GOLD, font_name="Arial")
    add_textbox(slide, Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.8), 
                title_text, font_size=32, bold=True, color=COLOR_DARK, font_name="Georgia")

# Helper to draw a thin gold horizontal separator line
def add_separator_line(slide, top=Inches(1.8)):
    # Create line using shapes
    slide.shapes.add_shape(
        1, # Rectangle
        Inches(1.0), top, Inches(11.333), Inches(0.02)
    ).fill.solid()
    # Recolor shape to Gold
    slide.shapes[-1].fill.fore_color.rgb = COLOR_GOLD
    slide.shapes[-1].line.color.rgb = COLOR_GOLD

# ==========================================================
# Slide 1: Welcome Page (Dark Theme)
# ==========================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, COLOR_DARK)
add_textbox(slide1, Inches(1), Inches(2.2), Inches(11.3), Inches(0.5), 
            "AUTOMATED TESTING FRAMEWORK", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(2.9), Inches(11.3), Inches(1.5), 
            "EventSphere Quality Assurance Report", font_size=42, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(4.7), Inches(11.3), Inches(1.0), 
            "An analysis of Cypress E2E, Selenium Webdriver, and Selenium IDE implementation", font_size=15, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# ==========================================================
# Slide 2: Agenda (Light Theme)
# ==========================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, COLOR_LIGHT)
add_slide_header(slide2, "Platform Overview", "Testing Agenda & Methodologies")
add_separator_line(slide2)

add_textbox(slide2, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.5), 
            "We have implemented three independent, automated testing layers to guarantee absolute stability, responsiveness, and performance of the EventSphere premium event management portal:", font_size=15, color=COLOR_MUTED)

add_bullets(slide2, Inches(1.5), Inches(3.2), Inches(10.3), Inches(3.5), [
    "Section 1: Selenium IDE Automated Tests - Visual codeless browser session recordings.",
    "Section 2: Selenium Webdriver Node.js Execution - Programmatic terminal-driven Chrome automation.",
    "Section 3: Cypress E2E Testing Suite - Complete client-side application specs, form submission, and route validation.",
    "Summary & Tool Comparison - Assessing stability, speed, and capabilities of each quality assurance strategy."
], font_size=16, color=COLOR_DARK)

# ==========================================================
# Slide 3: Section 1 Divider (Dark Theme)
# ==========================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, COLOR_DARK)
add_textbox(slide3, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5), 
            "SECTION 01", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide3, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.5), 
            "Selenium IDE Automated Tests", font_size=40, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide3, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Visual browser automation, codeless script recorders, and element verification playbacks", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# ==========================================================
# Slide 4: Section 1 Core Concepts (Light Theme)
# ==========================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, COLOR_LIGHT)
add_slide_header(slide4, "Section 1: Selenium IDE", "Codeless Visual Testing & Recording")
add_separator_line(slide4)

add_textbox(slide4, Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.5), 
            "Selenium IDE provides a rapid prototype automation layer by recording human actions in Chrome and generating assertions without writing code.\n\nIt is utilized for sanity testing critical user flows instantly when migrating codebases.", font_size=15, color=COLOR_MUTED)

add_bullets(slide4, Inches(6.5), Inches(2.2), Inches(5.8), Inches(4.5), [
    "Quick Test Initialization: Configures target site URL (Render host) directly in the extension dashboard.",
    "Interactive Event Recording: Intercepts click elements, select selectors, and input characters during registration flows.",
    "Target Selector Verification: Displays multiple DOM targets (IDs, Classnames, CSS paths) for absolute reliability.",
    "Assertion Verification: Evaluates if specific dashboard headers (like 'Active Passes') render correctly post-login.",
    "Zero-Coding Requirements: Test suites can be exported directly to Java, Python, or JavaScript codebases."
], font_size=14, color=COLOR_DARK)

# ==========================================================
# Slide 5: Section 2 Divider (Dark Theme)
# ==========================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5, COLOR_DARK)
add_textbox(slide5, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5), 
            "SECTION 02", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide5, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.5), 
            "Selenium Webdriver Terminal Scripts", font_size=40, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide5, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Programmatic Node.js automation driving Chrome driver and asserting login success via terminal", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# ==========================================================
# Slide 6: Section 2 Execution (Light Theme)
# ==========================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6, COLOR_LIGHT)
add_slide_header(slide6, "Section 2: Selenium Webdriver", "Node.js Programmatic Browser Automation")
add_separator_line(slide6)

add_textbox(slide6, Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.5), 
            "Our Selenium WebDriver script runs inside Node.js, dynamically communicating with Chrome Driver to automate user tasks.\n\nIt is executed via command-line and integrated into standard CI/CD deployment pipelines.", font_size=15, color=COLOR_MUTED)

add_bullets(slide6, Inches(6.5), Inches(2.2), Inches(5.8), Inches(4.5), [
    "JavaScript Integration: Uses 'selenium-webdriver' package directly in Node.js to configure Chrome browser sessions.",
    "Automated Input & Navigation: Instructs Chrome to open the login page, locate email/password inputs, and enter credentials.",
    "Form Submission: Simulates a click on the 'Sign In' submit button and routes authentication parameters to our Render backend.",
    "Dynamic Assertions: Waits for the URL to change to '/dashboard' to confirm the login workflow connected successfully.",
    "Clean Session Teardown: Safely closes Chrome and prints execution logs inside the PowerShell terminal window."
], font_size=14, color=COLOR_DARK)

# ==========================================================
# Slide 7: Section 3 Divider (Dark Theme)
# ==========================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7, COLOR_DARK)
add_textbox(slide7, Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.5), 
            "SECTION 03", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide7, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.5), 
            "Cypress End-to-End Suite", font_size=40, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide7, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Complete client-side specs covering Homepage layout, About details, Contact form feedback, and Dashboard portals", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# ==========================================================
# Slide 8: Section 3 Details (Light Theme)
# ==========================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide8, COLOR_LIGHT)
add_slide_header(slide8, "Section 3: Cypress E2E", "Cypress Client-Side Framework")
add_separator_line(slide8)

add_textbox(slide8, Inches(1.0), Inches(2.2), Inches(5.0), Inches(4.5), 
            "Cypress runs directly inside the browser environment, offering blazing fast spec execution, automatic waiting, and built-in screenshot debuggers.\n\nIt forms the core testing layer of EventSphere.", font_size=15, color=COLOR_MUTED)

add_bullets(slide8, Inches(6.5), Inches(2.2), Inches(5.8), Inches(4.5), [
    "Homepage Check: Verifies luxury brand header ('EVENTSPHERE') and navigation links exist and are fully interactive.",
    "About Page Verification: Visits '/about' and asserts that core text paragraphs load successfully.",
    "Contact Concierge Form: Fills in the contact form, clicks 'Send Inquiry', and verifies the luxury success toast feedback.",
    "Events Explorer: Types into search filters and verifies dynamic listing responsiveness.",
    "Redux Store & Router Checks: Logs in and navigates dashboard tabs (My Tickets, Wishlist, Settings) to verify client routes."
], font_size=14, color=COLOR_DARK)

# ==========================================================
# Slide 9: Summary & Comparison (Light Theme)
# ==========================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide9, COLOR_LIGHT)
add_slide_header(slide9, "Summary", "Automated Testing Tools Comparison")
add_separator_line(slide9)

add_textbox(slide9, Inches(1.0), Inches(2.2), Inches(3.6), Inches(4.5), 
            "Selenium IDE\n\n• Codeless recorder\n• Best for rapid sanity prototypes\n• Fast configurations\n• Requires browser extension", font_size=14, bold=True, color=COLOR_DARK)

add_textbox(slide9, Inches(4.8), Inches(2.2), Inches(3.6), Inches(4.5), 
            "Selenium WebDriver\n\n• Code-based (Node.js)\n• Best for CI/CD integration\n• Runs headless in shell\n• Excellent cross-browser support", font_size=14, bold=True, color=COLOR_DARK)

add_textbox(slide9, Inches(8.6), Inches(2.2), Inches(3.6), Inches(4.5), 
            "Cypress E2E\n\n• Direct browser architecture\n• Super fast execution times\n• Built-in time travel & debugging\n• Comprehensive UI/UX tests", font_size=14, bold=True, color=COLOR_DARK)

# ==========================================================
# Slide 10: Conclusion (Dark Theme)
# ==========================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide10, COLOR_DARK)
add_textbox(slide10, Inches(1), Inches(2.2), Inches(11.3), Inches(0.5), 
            "QUALITY ASSURANCE SUCCESS", font_size=12, bold=True, color=COLOR_GOLD, font_name="Arial", align=PP_ALIGN.CENTER)
add_textbox(slide10, Inches(1), Inches(2.9), Inches(11.3), Inches(1.5), 
            "All Test Suites Verified & Operational", font_size=38, bold=True, color=COLOR_WHITE, font_name="Georgia", align=PP_ALIGN.CENTER)
add_textbox(slide10, Inches(1), Inches(4.5), Inches(11.3), Inches(1.0), 
            "Cypress and Selenium Webdriver scripts compile and execute successfully with zero errors.", font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# Save presentation
output_pptx = r"c:\Users\Yash Bhatt\Desktop\Demo\EventSphere_Testing_Overview.pptx"
prs.save(output_pptx)
print(f"PowerPoint Presentation saved successfully: {output_pptx}")
